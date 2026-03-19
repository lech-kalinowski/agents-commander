#!/usr/bin/env python3
"""Generate Commander Protocol PowerPoint presentation with architecture diagram.
VC-pitch ready for Silicon Valley investors."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──
BG_DARK    = RGBColor(0x0D, 0x11, 0x17)
BG_CARD    = RGBColor(0x16, 0x1B, 0x22)
CYAN       = RGBColor(0x00, 0xD4, 0xFF)
CYAN_DIM   = RGBColor(0x00, 0x8B, 0xAA)
GREEN      = RGBColor(0x3F, 0xB9, 0x50)
YELLOW     = RGBColor(0xFF, 0xD6, 0x00)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)
RED        = RGBColor(0xE0, 0x40, 0x40)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x8B, 0x94, 0x9E)
LIGHT_GRAY = RGBColor(0xC9, 0xD1, 0xD9)
BLUE       = RGBColor(0x58, 0xA6, 0xFF)
PURPLE     = RGBColor(0xBC, 0x8C, 0xFF)
GEMINI_BLUE = RGBColor(0x44, 0x85, 0xF4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, w, h, fill_color, border_color=None, border_w=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_w
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def set_text(shape, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.paragraphs[0].alignment = align
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return tf

def add_text_box(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return tf

def add_bullet_slide(slide, items, left, top, w, size=18, color=LIGHT_GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, w, Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = spacing
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def add_arrow(slide, x1, y1, x2, y2, color=CYAN, width=Pt(2.5)):
    """Add a line with arrow from (x1,y1) to (x2,y2)."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        x1, y1, x2, y2
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead at the end
    connector.end_style = 'triangle'  # may not work, set via XML
    from pptx.oxml.ns import qn
    ln = connector.line._ln
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)
    return connector


# ===============================================================
# SLIDE 1 -- Title
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, BG_DARK)

# Accent line
add_rect(slide, Inches(4), Inches(1.5), Inches(5.3), Pt(3), CYAN)

add_text_box(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
             "AGENTS COMMANDER", size=48, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

# Version badge
badge = add_shape(slide, Inches(5.9), Inches(2.9), Inches(1.5), Inches(0.45), CYAN, CYAN, Pt(1))
set_text(badge, "v0.1.0", size=16, color=BG_DARK, bold=True)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
             "Commander Protocol", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.8),
             "A novel inter-agent communication protocol for terminal-based AI agents",
             size=20, color=GRAY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.6),
             "Stop copy-pasting between AI agents. Let them work together.",
             size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Accent line bottom
add_rect(slide, Inches(4), Inches(5.8), Inches(5.3), Pt(3), CYAN)

add_text_box(slide, Inches(1), Inches(6.1), Inches(11.3), Inches(0.6),
             "Lech Kalinowski  |  agents-commander.dev  |  2026", size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ===============================================================
# SLIDE 2 -- Product Screenshot
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "PRODUCT SCREENSHOT", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(3), CYAN)

# Large screenshot placeholder
placeholder = add_shape(slide, Inches(1.2), Inches(1.7), Inches(10.9), Inches(4.6),
                        RGBColor(0x10, 0x15, 0x1C), CYAN_DIM, Pt(2))
set_text(placeholder,
         "Screenshot Placeholder\n\nReplace with actual product screenshot\n\n10.9\" x 4.6\" (1920x810 px recommended)",
         size=20, color=GRAY)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.6),
             "Live demo: agents-commander.dev",
             size=16, color=CYAN_DIM, align=PP_ALIGN.CENTER)


# ===============================================================
# SLIDE 3 -- Market Opportunity
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "MARKET OPPORTUNITY", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(3), CYAN)

# Market stats cards
market_stats = [
    ("$24B", "AI Developer Tools\nMarket by 2028", CYAN),
    ("2M+", "Developers Using\nAI Coding Agents", GREEN),
    ("0", "Multi-Agent Orchestration\nStandards Exist", YELLOW),
]

for i, (big_num, desc, color) in enumerate(market_stats):
    x = Inches(0.6 + i * 4.2)
    card = add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(3.2), BG_CARD, color, Pt(2))

    # Big number
    add_text_box(slide, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(1.0),
                 big_num, size=52, color=color, bold=True, align=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Inches(0.3), Inches(3.3), Inches(3.2), Inches(1.2),
                 desc, size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom statement
add_rect(slide, Inches(2.5), Inches(5.5), Inches(8.3), Pt(2), CYAN_DIM)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.8),
             "Commander Protocol is the missing infrastructure layer",
             size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ===============================================================
# SLIDE 4 -- Problem
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "THE PROBLEM", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Pt(3), CYAN)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
             "AI coding agents (Claude Code, Codex CLI, Gemini CLI) are powerful individually,\n"
             "but they operate in isolated terminal sessions with no way to collaborate.",
             size=20, color=LIGHT_GRAY)

# Problem cards
problems = [
    ("No Standard Protocol", "No industry standard for CLI agents\nto communicate with each other", RED),
    ("Isolated Processes", "Each agent runs in its own PTY\nwith no shared state or messaging", ORANGE),
    ("API-Only Solutions", "Existing frameworks (MCP, LangGraph)\nrequire API integration, not terminal", YELLOW),
]

for i, (title, desc, color) in enumerate(problems):
    x = Inches(0.8 + i * 4.1)
    card = add_shape(slide, x, Inches(3.2), Inches(3.7), Inches(2.8), BG_CARD, color, Pt(2))
    add_text_box(slide, x + Inches(0.3), Inches(3.5), Inches(3.1), Inches(0.6),
                 title, size=22, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(4.2), Inches(3.1), Inches(1.5),
                 desc, size=16, color=LIGHT_GRAY)


# ===============================================================
# SLIDE 5 -- Solution Overview
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "THE SOLUTION: COMMANDER PROTOCOL", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(5), Pt(3), CYAN)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1),
             "A sideband text-marker protocol that enables inter-agent communication\n"
             "by scanning terminal output and injecting input \u2014 zero API integration required.",
             size=20, color=LIGHT_GRAY)

# Protocol block example
code_bg = add_shape(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(3.2), RGBColor(0x0D, 0x11, 0x17), CYAN_DIM, Pt(1.5))

code_lines = [
    ("===COMMANDER:SEND:codex:2===",  CYAN),
    ("",                               WHITE),
    ("Review the file src/app.ts and fix",  WHITE),
    ("any TypeScript errors you find.",      WHITE),
    ("",                               WHITE),
    ("===COMMANDER:END===",            CYAN),
]

for i, (line, color) in enumerate(code_lines):
    add_text_box(slide, Inches(2.0), Inches(3.2 + i * 0.4), Inches(9), Inches(0.4),
                 line, size=18, color=color, bold=(color == CYAN))

add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
             "Agent outputs this marker block  \u2192  Commander detects it  \u2192  Routes to Codex in Panel 2",
             size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ===============================================================
# SLIDE 6 -- Architecture Diagram
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
             "ARCHITECTURE", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(0.95), Inches(3), Pt(3), CYAN)

# ── Central Orchestrator ──
orch_x, orch_y = Inches(4.5), Inches(1.6)
orch_w, orch_h = Inches(4.3), Inches(1.2)
orch = add_shape(slide, orch_x, orch_y, orch_w, orch_h, RGBColor(0x1A, 0x2A, 0x3A), CYAN, Pt(3))
set_text(orch, "ORCHESTRATOR\nMessage Router  |  Task Queue  |  Protocol Injector", size=14, color=CYAN, bold=True)

# ── Panel boxes ──
panel_data = [
    ("Panel 1",  "Claude Code",   "Anthropic",  Inches(0.5),  GREEN),
    ("Panel 2",  "Codex CLI",     "OpenAI",     Inches(4.5),  BLUE),
    ("Panel 3",  "Gemini CLI",    "Google",     Inches(8.5),  GEMINI_BLUE),
]

panel_tops = Inches(4.0)
panel_w = Inches(3.6)
panel_h = Inches(2.8)

for (label, agent, vendor, px, color) in panel_data:
    # Panel frame
    panel = add_shape(slide, px, panel_tops, panel_w, panel_h, BG_CARD, color, Pt(2))

    # Panel header bar
    hdr = add_rect(slide, px + Inches(0.05), panel_tops + Inches(0.05), panel_w - Inches(0.1), Inches(0.5), color)
    set_text(hdr, f"{label}  \u2014  {agent}", size=14, color=BG_DARK, bold=True)

    # Agent info
    add_text_box(slide, px + Inches(0.2), panel_tops + Inches(0.7), panel_w - Inches(0.4), Inches(0.4),
                 vendor, size=13, color=GRAY)

    # Scanner box
    scan = add_shape(slide, px + Inches(0.15), panel_tops + Inches(1.1), panel_w - Inches(0.3), Inches(0.7),
                     RGBColor(0x10, 0x15, 0x1C), CYAN_DIM, Pt(1))
    set_text(scan, "ProtocolScanner\nANSI strip \u2192 buffer \u2192 regex match", size=11, color=CYAN_DIM)

    # PTY box
    pty = add_shape(slide, px + Inches(0.15), panel_tops + Inches(1.95), panel_w - Inches(0.3), Inches(0.65),
                    RGBColor(0x10, 0x15, 0x1C), GRAY, Pt(1))
    set_text(pty, "PTY Process (stdin/stdout)", size=11, color=GRAY)

# ── Arrows: Orchestrator <-> Panels ──
# Panel 1 (left)
add_arrow(slide, Inches(2.3), panel_tops, Inches(4.5), orch_y + orch_h, GREEN)
add_arrow(slide, Inches(5.0), orch_y + orch_h, Inches(2.8), panel_tops, GREEN)

# Panel 2 (center)
add_arrow(slide, Inches(6.3), panel_tops, Inches(6.3), orch_y + orch_h, BLUE)
add_arrow(slide, Inches(7.0), orch_y + orch_h, Inches(7.0), panel_tops, BLUE)

# Panel 3 (right)
add_arrow(slide, Inches(10.3), panel_tops, Inches(8.8), orch_y + orch_h, GEMINI_BLUE)
add_arrow(slide, Inches(8.3), orch_y + orch_h, Inches(9.8), panel_tops, GEMINI_BLUE)

# Arrow labels
add_text_box(slide, Inches(0.5), Inches(3.3), Inches(2.5), Inches(0.5),
             "\u2191 detected markers", size=11, color=GREEN, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.2), Inches(3.3), Inches(2.5), Inches(0.5),
             "\u2193 inject tasks", size=11, color=BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.5), Inches(3.3), Inches(2.5), Inches(0.5),
             "\u2191 scan output", size=11, color=GEMINI_BLUE, align=PP_ALIGN.CENTER)


# ===============================================================
# SLIDE 7 -- Message Flow
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "MESSAGE FLOW", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Pt(3), CYAN)

steps = [
    ("1", "INJECT",   "Ctrl+P injects protocol instructions\ninto agent's stdin",                    CYAN),
    ("2", "OUTPUT",   "Agent includes COMMANDER markers\nin its terminal output",                     GREEN),
    ("3", "SCAN",     "ProtocolScanner strips ANSI, buffers\nlines, matches regex markers",           BLUE),
    ("4", "ROUTE",    "Orchestrator queues message,\nlaunches target agent if needed",                PURPLE),
    ("5", "DELIVER",  "Task text sent to target agent's\nstdin in chunks + Enter",                    YELLOW),
    ("6", "RESPOND",  "Target agent processes task,\ncan reply back via same protocol",               GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = Inches(1.7 + i * 0.9)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.shadow.inherit = False
    set_text(circle, num, size=18, color=BG_DARK, bold=True)

    # Title
    add_text_box(slide, Inches(1.9), y + Inches(0.02), Inches(2.2), Inches(0.5),
                 title, size=20, color=color, bold=True)

    # Description
    add_text_box(slide, Inches(4.2), y + Inches(0.02), Inches(8), Inches(0.6),
                 desc, size=16, color=LIGHT_GRAY)

    # Connector line to next
    if i < len(steps) - 1:
        add_rect(slide, Inches(1.27), y + Inches(0.6), Pt(2), Inches(0.3), color)


# ===============================================================
# SLIDE 8 -- Key Components
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "KEY COMPONENTS", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Pt(3), CYAN)

components = [
    ("ProtocolScanner",       "Stateful line-by-line scanner attached to each\n"
                              "terminal panel. Strips ANSI codes, buffers text,\n"
                              "detects COMMANDER markers via regex.",                CYAN),
    ("Orchestrator",          "Central message router with per-panel task queues.\n"
                              "Launches agents on demand, injects protocol,\n"
                              "sends text in chunks to avoid buffer overflow.",      GREEN),
    ("Commander Protocol",    "Text-based marker format embedded in agent output.\n"
                              "===COMMANDER:SEND:agent:panel===\n"
                              "===COMMANDER:END===",                                 YELLOW),
    ("Terminal Panel (PTY)",  "Managed child process with stdin/stdout pipes.\n"
                              "VTerm emulator, ANSI processing, mouse/key\n"
                              "forwarding, 2000-line scrollback.",                   BLUE),
]

for i, (title, desc, color) in enumerate(components):
    x = Inches(0.6 + (i % 2) * 6.3)
    y = Inches(1.6 + (i // 2) * 2.7)
    card = add_shape(slide, x, y, Inches(5.9), Inches(2.3), BG_CARD, color, Pt(2))
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.3), Inches(0.5),
                 title, size=22, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(5.3), Inches(1.3),
                 desc, size=15, color=LIGHT_GRAY)


# ===============================================================
# SLIDE 9 -- What Makes It Unique
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "WHAT MAKES IT UNIQUE", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(3), CYAN)

comparisons = [
    ("Commander Protocol",    ["Terminal-native (PTY scanning)",
                               "Works with ANY CLI agent",
                               "No code changes to agents",
                               "Text markers in stdout",
                               "Sideband approach"],           CYAN),
    ("MCP (Anthropic)",       ["API-level (JSON-RPC)",
                               "Requires SDK integration",
                               "Server/client architecture",
                               "Structured tool calls",
                               "Deep integration needed"],     PURPLE),
    ("LangGraph / CrewAI",    ["Python framework",
                               "Programmatic orchestration",
                               "Requires custom agent code",
                               "API-based communication",
                               "Framework lock-in"],           ORANGE),
]

for i, (title, items, color) in enumerate(comparisons):
    x = Inches(0.6 + i * 4.2)
    card = add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(4.8), BG_CARD, color, Pt(2))

    # Header
    hdr = add_rect(slide, x + Inches(0.05), Inches(1.85), Inches(3.7), Inches(0.6), color)
    set_text(hdr, title, size=18, color=BG_DARK, bold=True)

    for j, item in enumerate(items):
        prefix = "\u2713 " if i == 0 else "\u2022 "
        item_color = GREEN if i == 0 else LIGHT_GRAY
        add_text_box(slide, x + Inches(0.3), Inches(2.7 + j * 0.55), Inches(3.2), Inches(0.5),
                     prefix + item, size=15, color=item_color)


# ===============================================================
# SLIDE 10 -- Beyond Agents Commander
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "BEYOND AGENTS COMMANDER", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4.5), Pt(3), CYAN)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.8),
             "The protocol is portable \u2014 if an agent can print text and read text,\n"
             "it can participate. No SDK, no API keys, no integration work.",
             size=20, color=LIGHT_GRAY)

use_cases = [
    ("Shell Pipe Relay",
     "Daemon watches tmux/screen panes,\nscans output, routes messages\nbetween existing terminal sessions",
     CYAN),
    ("CI/CD Coordination",
     "AI agents in CI pipeline jobs\nexchange tasks via log output.\nAgent A reviews, Agent B fixes",
     GREEN),
    ("VS Code Extension",
     "IDE runs agent CLIs in terminal\ntabs, extension scans output\nand relays between them",
     BLUE),
    ("Docker Sidecar",
     "Each agent in its own container.\nSidecar attaches to stdout\nand routes via shared volume",
     PURPLE),
    ("MCP Bridge",
     "Commander-to-MCP adapter:\ntranslates markers into MCP\ntool calls and vice versa",
     YELLOW),
    ("npm Package",
     "Standalone ProtocolScanner library.\nnpm install commander-protocol\nAdd to any Node.js app",
     ORANGE),
]

for i, (title, desc, color) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(2.8 + row * 2.3)
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.0), BG_CARD, color, Pt(2))
    add_text_box(slide, x + Inches(0.25), y + Inches(0.15), Inches(3.3), Inches(0.4),
                 title, size=18, color=color, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.6), Inches(3.3), Inches(1.2),
                 desc, size=13, color=LIGHT_GRAY)


# ===============================================================
# SLIDE 11 -- Traction & Roadmap
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "TRACTION & ROADMAP", size=36, color=CYAN, bold=True)
add_rect(slide, Inches(0.8), Inches(1.1), Inches(4), Pt(3), CYAN)

# Left card: v0.1.0 Launch
left_card = add_shape(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(5.0), BG_CARD, GREEN, Pt(2))
hdr_left = add_rect(slide, Inches(0.65), Inches(1.75), Inches(5.7), Inches(0.7), GREEN)
set_text(hdr_left, "v0.1.0 \u2014 Launch", size=22, color=BG_DARK, bold=True)

launch_items = [
    "\u2713  10 agent types supported (Claude, Codex, Gemini, Copilot, Aider, ...)",
    "\u2713  120 built-in prompt templates",
    "\u2713  Commander Protocol v2 (SEND, REPLY, BROADCAST, STATUS, QUERY + ACK)",
    "\u2713  Dual-panel file manager with preview",
    "\u2713  Full terminal emulation per panel",
    "\u2713  Built-in Markdown editor with live preview",
    "\u2713  Agent skill file discovery",
    "\u2713  Dark theme with customizable layouts",
]

for i, item in enumerate(launch_items):
    add_text_box(slide, Inches(1.0), Inches(2.7 + i * 0.45), Inches(5.0), Inches(0.45),
                 item, size=14, color=GREEN)

# Right card: What's Next
right_card = add_shape(slide, Inches(6.9), Inches(1.7), Inches(5.8), Inches(5.0), BG_CARD, CYAN, Pt(2))
hdr_right = add_rect(slide, Inches(6.95), Inches(1.75), Inches(5.7), Inches(0.7), CYAN)
set_text(hdr_right, "What's Next", size=22, color=BG_DARK, bold=True)

roadmap_items = [
    "\u2192  npm protocol library (standalone scanner)",
    "\u2192  VS Code extension for IDE integration",
    "\u2192  CI/CD pipeline integration",
    "\u2192  Session persistence and replay",
    "\u2192  Plugin system for custom agents",
    "\u2192  MCP bridge adapter",
    "\u2192  Multi-workspace support",
    "\u2192  Team collaboration features",
]

for i, item in enumerate(roadmap_items):
    add_text_box(slide, Inches(7.3), Inches(2.7 + i * 0.45), Inches(5.0), Inches(0.45),
                 item, size=14, color=CYAN)


# ===============================================================
# SLIDE 12 -- Closing / Summary
# ===============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

# Top accent line
add_rect(slide, Inches(3.5), Inches(1.0), Inches(6.3), Pt(3), CYAN)

add_text_box(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(1.2),
             "AGENTS COMMANDER", size=52, color=CYAN, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(0.8),
             "The infrastructure layer for multi-agent AI collaboration",
             size=26, color=WHITE, bold=False, align=PP_ALIGN.CENTER)

# Key stats bar
stats_bar = add_shape(slide, Inches(2.0), Inches(3.5), Inches(9.3), Inches(0.8), BG_CARD, CYAN_DIM, Pt(1.5))
set_text(stats_bar, "10 agent types  \u2022  120 templates  \u2022  5 protocol commands  \u2022  Zero integration",
         size=22, color=LIGHT_GRAY, bold=True)

# Install command
install_bg = add_shape(slide, Inches(3.0), Inches(4.7), Inches(7.3), Inches(0.8),
                       RGBColor(0x0D, 0x11, 0x17), GREEN, Pt(2))
set_text(install_bg, "npm install -g agents-commander", size=22, color=GREEN, bold=True)

add_text_box(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.4),
             "Try it now:", size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Bottom accent line
add_rect(slide, Inches(3.5), Inches(5.9), Inches(6.3), Pt(3), CYAN)

add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
             "Lech Kalinowski  |  agents-commander.dev  |  CC BY-NC 4.0",
             size=16, color=GRAY, align=PP_ALIGN.CENTER)


# ── Save ──
output = "/Users/lech.kalinowski/products/agents commander/Agents_Commander_Protocol.pptx"
prs.save(output)
print(f"Saved: {output}")
