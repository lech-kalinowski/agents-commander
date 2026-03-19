#!/usr/bin/env python3
"""Build the Agents Commander Protocol presentation with green military theme."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme colors ──
BG_BLACK = RGBColor(0x04, 0x0A, 0x04)
BG_DARK = RGBColor(0x06, 0x12, 0x06)
BG_CARD = RGBColor(0x0A, 0x18, 0x0A)
GREEN_BRIGHT = RGBColor(0x00, 0xFF, 0x41)
GREEN_MID = RGBColor(0x00, 0xCC, 0x33)
GREEN_DIM = RGBColor(0x00, 0x88, 0x22)
GREEN_MUTED = RGBColor(0x4A, 0x6A, 0x4A)
TEXT_PRIMARY = RGBColor(0xD0, 0xE8, 0xD0)
TEXT_SECONDARY = RGBColor(0x7A, 0x9A, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xFF, 0x44, 0x44)
YELLOW = RGBColor(0x88, 0xCC, 0x44)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
LOGO_PATH = os.path.join(PROJECT_DIR, "assets", "logo.png")
OUTPUT_PATH = os.path.join(PROJECT_DIR, "Agents_Commander_Protocol.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_slide_bg(slide, color=BG_BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.rotation = 0
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_PRIMARY,
             bold=False, alignment=PP_ALIGN.LEFT, font_name='Consolas'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16, font_name='Consolas'):
    """lines is a list of (text, color, bold) tuples."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, color, bold) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


def slide_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
             f"{num}/{total}", font_size=11, color=GREEN_MUTED,
             alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12

# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.167), Inches(0.6), Inches(3.0))

add_text(slide, Inches(0), Inches(3.7), W, Inches(1),
         "AGENTS COMMANDER", font_size=54, color=GREEN_BRIGHT,
         bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(4.7), W, Inches(0.6),
         "Multi-Panel Terminal for AI Agent Collaboration", font_size=24,
         color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(5.6), W, Inches(0.5),
         "The Commander Protocol  \u2022  Inter-Agent Communication  \u2022  Terminal-Native",
         font_size=16, color=GREEN_MUTED, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(6.5), W, Inches(0.4),
         "by Lech Kalinowski  \u2022  CC BY-NC 4.0  \u2022  v0.1.0",
         font_size=13, color=GREEN_MUTED, alignment=PP_ALIGN.CENTER)

slide_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2 — The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "THE PROBLEM", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "AI agents are powerful. But isolated.", font_size=38, color=GREEN_BRIGHT, bold=True)

problems = [
    ("No Standard Protocol", "No universal way for terminal-based AI agents to communicate.\nEach tool invents its own integration, creating fragmentation."),
    ("Isolated Processes", "Claude cannot talk to Codex. Gemini cannot hand off to Aider.\nEach agent operates in a vacuum, blind to other agents."),
    ("API-Only Solutions", "Existing multi-agent frameworks require API keys, custom code,\nand cloud infrastructure. Nothing works natively in the terminal."),
]

for i, (title, desc) in enumerate(problems):
    left = Inches(0.8 + i * 4.0)
    top = Inches(2.5)
    add_rect(slide, left, top, Inches(3.6), Inches(3.5), BG_CARD, GREEN_DIM)
    add_text(slide, left + Inches(0.3), top + Inches(0.4), Inches(3.0), Inches(0.5),
             "\u2716", font_size=28, color=RED, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.0), Inches(3.0), Inches(0.5),
             title, font_size=20, color=TEXT_PRIMARY, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.7), Inches(3.0), Inches(1.6),
             desc, font_size=14, color=TEXT_SECONDARY)

slide_number(slide, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3 — The Solution Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "THE SOLUTION", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "One terminal. Multiple agents. Zero integration.", font_size=38,
         color=GREEN_BRIGHT, bold=True)

add_text(slide, Inches(0.8), Inches(2.2), Inches(11), Inches(0.6),
         "The Commander Protocol is a text-based messaging format that any agent can produce and any terminal can route.",
         font_size=18, color=TEXT_SECONDARY)

steps = [
    ("1", "INJECT", "Skill Files", "Commander injects a skill file into each agent,\nteaching it the protocol message format.\nNo API integration needed."),
    ("2", "SCAN", "Output Monitoring", "Commander continuously scans each agent's\nterminal output for protocol markers.\nPure text matching \u2014 works with any agent."),
    ("3", "ROUTE", "Message Delivery", "When a marker is detected, Commander extracts\nthe message and injects it into the target\nagent's stdin as a new prompt."),
]

for i, (num, keyword, title, desc) in enumerate(steps):
    left = Inches(0.8 + i * 4.0)
    top = Inches(3.3)
    add_rect(slide, left, top, Inches(3.6), Inches(3.5), BG_CARD, GREEN_DIM)
    # Step number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN_DIM
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = GREEN_BRIGHT
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Consolas'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text(slide, left + Inches(1.0), top + Inches(0.35), Inches(2.3), Inches(0.4),
             f"{keyword} \u2014 {title}", font_size=16, color=GREEN_BRIGHT, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.2), Inches(3.0), Inches(2.0),
             desc, font_size=14, color=TEXT_SECONDARY)

slide_number(slide, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4 — Protocol Commands
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "COMMANDER PROTOCOL", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "5 commands. One end marker. Fully bidirectional.", font_size=38,
         color=GREEN_BRIGHT, bold=True)

commands = [
    ("SEND", "Direct message to a specific agent", "===COMMANDER:SEND:codex:2===\nWrite tests for src/auth/\n===COMMANDER:END==="),
    ("REPLY", "Respond to whoever last messaged you", "===COMMANDER:REPLY===\n12 tests passing, 0 failing.\n===COMMANDER:END==="),
    ("BROADCAST", "Send to all connected agents", "===COMMANDER:BROADCAST===\nPhase 1 complete. Begin phase 2.\n===COMMANDER:END==="),
    ("STATUS", "Progress toast in UI (not sent to agents)", "===COMMANDER:STATUS===\nAnalyzing file 5 of 10...\n===COMMANDER:END==="),
    ("QUERY", "Ask Commander what agents are running", "===COMMANDER:QUERY===\nagents\n===COMMANDER:END==="),
]

for i, (cmd, desc, example) in enumerate(commands):
    left = Inches(0.8)
    top = Inches(2.2 + i * 1.0)
    # Command name
    add_text(slide, left, top, Inches(1.5), Inches(0.4),
             cmd, font_size=18, color=GREEN_BRIGHT, bold=True)
    # Description
    add_text(slide, left + Inches(1.7), top, Inches(4.5), Inches(0.4),
             desc, font_size=15, color=TEXT_SECONDARY)

# Code example box on right
add_rect(slide, Inches(7.5), Inches(2.0), Inches(5.2), Inches(5.0), BG_CARD, GREEN_DIM)
add_text(slide, Inches(7.7), Inches(2.1), Inches(4.8), Inches(0.4),
         "Protocol Examples", font_size=12, color=GREEN_MUTED, bold=True)

code_lines = [
    ("// SEND \u2014 direct message", GREEN_MUTED, False),
    ("===COMMANDER:SEND:codex:2===", GREEN_BRIGHT, True),
    ("Write unit tests for src/auth/", TEXT_PRIMARY, False),
    ("===COMMANDER:END===", GREEN_BRIGHT, True),
    ("", TEXT_PRIMARY, False),
    ("// REPLY \u2014 respond to sender", GREEN_MUTED, False),
    ("===COMMANDER:REPLY===", GREEN_BRIGHT, True),
    ("12 tests passing, 0 failing.", TEXT_PRIMARY, False),
    ("===COMMANDER:END===", GREEN_BRIGHT, True),
    ("", TEXT_PRIMARY, False),
    ("// BROADCAST \u2014 all agents", GREEN_MUTED, False),
    ("===COMMANDER:BROADCAST===", GREEN_BRIGHT, True),
    ("Phase 1 complete. Begin phase 2.", TEXT_PRIMARY, False),
    ("===COMMANDER:END===", GREEN_BRIGHT, True),
    ("", TEXT_PRIMARY, False),
    ("\u2192 Sender gets ACK after delivery", GREEN_MID, False),
]

add_multiline(slide, Inches(7.7), Inches(2.5), Inches(4.8), Inches(4.3),
              code_lines, font_size=13)

slide_number(slide, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5 — Message Format Detail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "MESSAGE FORMAT", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "Pure text markers. Any agent that can print can use them.", font_size=38,
         color=GREEN_BRIGHT, bold=True)

# Format breakdown box
add_rect(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.2), BG_CARD, GREEN_DIM)

fmt_lines = [
    ("===COMMANDER", GREEN_BRIGHT, True),
    (":SEND", YELLOW, True),
    (":codex", GREEN_MID, True),
    (":2", RGBColor(0x88, 0xCC, 0x44), True),
    ("===", GREEN_BRIGHT, True),
]
# Build as single text with explanation below
add_text(slide, Inches(2.0), Inches(2.8), Inches(9.0), Inches(0.6),
         "===COMMANDER:SEND:codex:2===", font_size=28, color=GREEN_BRIGHT, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.0), Inches(3.5), Inches(9.0), Inches(0.5),
         "      marker       action     target agent   panel number      marker",
         font_size=14, color=GREEN_MUTED, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.0), Inches(4.0), Inches(9.0), Inches(0.5),
         "{message content \u2014 any text, any length, multiple lines}",
         font_size=18, color=TEXT_PRIMARY, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.0), Inches(4.4), Inches(9.0), Inches(0.5),
         "===COMMANDER:END===", font_size=28, color=GREEN_BRIGHT, bold=True,
         alignment=PP_ALIGN.CENTER)

# Delivery flow
add_text(slide, Inches(0.8), Inches(5.3), Inches(12), Inches(0.4),
         "Delivery Flow:", font_size=16, color=GREEN_MID, bold=True)

flow_items = [
    "1. Agent prints protocol markers to stdout",
    "2. ProtocolScanner detects markers in real-time (ANSI-stripped, chunked)",
    "3. Orchestrator extracts command + payload, validates target",
    "4. Message injected into target agent's stdin via bracketed paste",
    "5. Sender receives ACK: [Commander] Message delivered to codex in Panel 2 (OK)",
]
add_multiline(slide, Inches(0.8), Inches(5.7), Inches(12), Inches(2.0),
              [(line, TEXT_SECONDARY, False) for line in flow_items], font_size=14)

slide_number(slide, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6 — Architecture: How It Works
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "ARCHITECTURE", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "See it in action", font_size=38, color=GREEN_BRIGHT, bold=True)

# Panel 1 — Claude
add_rect(slide, Inches(0.8), Inches(2.3), Inches(3.6), Inches(4.5), BG_CARD, GREEN_DIM)
add_text(slide, Inches(1.0), Inches(2.4), Inches(3.2), Inches(0.4),
         "\u25CF Claude Code          Panel 1", font_size=13, color=GREEN_MID, bold=True)
panel1_lines = [
    ("Analyzing codebase...", TEXT_SECONDARY, False),
    ("Found auth vulnerability.", TEXT_SECONDARY, False),
    ("", TEXT_PRIMARY, False),
    ("===COMMANDER:SEND:codex:2===", GREEN_BRIGHT, True),
    ("Fix the SQL injection in", TEXT_PRIMARY, False),
    ("src/auth/login.ts line 42", TEXT_PRIMARY, False),
    ("===COMMANDER:END===", GREEN_BRIGHT, True),
    ("", TEXT_PRIMARY, False),
    ("[Commander] Delivered to", GREEN_MID, False),
    ("codex in Panel 2 (OK)", GREEN_MID, False),
]
add_multiline(slide, Inches(1.0), Inches(2.9), Inches(3.2), Inches(3.8), panel1_lines, font_size=12)

# Arrow 1
add_text(slide, Inches(4.5), Inches(4.2), Inches(1.0), Inches(0.5),
         "\u2794", font_size=32, color=GREEN_BRIGHT, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(4.7), Inches(1.0), Inches(0.3),
         "detect", font_size=11, color=GREEN_MUTED, alignment=PP_ALIGN.CENTER)

# Commander center
add_rect(slide, Inches(5.3), Inches(3.2), Inches(2.7), Inches(2.5), BG_CARD, GREEN_BRIGHT, Pt(2))
add_text(slide, Inches(5.3), Inches(3.5), Inches(2.7), Inches(0.5),
         "\u2699", font_size=32, alignment=PP_ALIGN.CENTER, color=GREEN_BRIGHT)
add_text(slide, Inches(5.3), Inches(4.1), Inches(2.7), Inches(0.4),
         "COMMANDER", font_size=14, color=GREEN_BRIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.3), Inches(4.5), Inches(2.7), Inches(0.5),
         "scan \u2022 parse \u2022 route", font_size=12, color=GREEN_MUTED, alignment=PP_ALIGN.CENTER)

# Arrow 2
add_text(slide, Inches(8.1), Inches(4.2), Inches(1.0), Inches(0.5),
         "\u2794", font_size=32, color=GREEN_BRIGHT, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.1), Inches(4.7), Inches(1.0), Inches(0.3),
         "inject", font_size=11, color=GREEN_MUTED, alignment=PP_ALIGN.CENTER)

# Panel 2 — Codex
add_rect(slide, Inches(9.0), Inches(2.3), Inches(3.6), Inches(4.5), BG_CARD, GREEN_DIM)
add_text(slide, Inches(9.2), Inches(2.4), Inches(3.2), Inches(0.4),
         "\u25CF Codex CLI             Panel 2", font_size=13, color=GREEN_MID, bold=True)
panel2_lines = [
    ("[Commander: from Claude]", GREEN_MID, True),
    ("Fix the SQL injection in", TEXT_SECONDARY, False),
    ("src/auth/login.ts line 42", TEXT_SECONDARY, False),
    ("", TEXT_PRIMARY, False),
    ("On it. Applying parameterized", TEXT_SECONDARY, False),
    ("query pattern...", TEXT_SECONDARY, False),
    ("", TEXT_PRIMARY, False),
    ("===COMMANDER:REPLY===", GREEN_BRIGHT, True),
    ("Fixed. Tests passing.", TEXT_PRIMARY, False),
    ("===COMMANDER:END===", GREEN_BRIGHT, True),
]
add_multiline(slide, Inches(9.2), Inches(2.9), Inches(3.2), Inches(3.8), panel2_lines, font_size=12)

slide_number(slide, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7 — Protocol Implementation Details
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "IMPLEMENTATION", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "Protocol Scanner & Orchestrator", font_size=38, color=GREEN_BRIGHT, bold=True)

# Left column — Scanner
add_rect(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(4.8), BG_CARD, GREEN_DIM)
add_text(slide, Inches(1.1), Inches(2.5), Inches(5.2), Inches(0.4),
         "ProtocolScanner", font_size=18, color=GREEN_BRIGHT, bold=True)

scanner_lines = [
    ("\u2022 Real-time line-by-line scanning of agent output", TEXT_SECONDARY, False),
    ("\u2022 ANSI escape codes stripped before matching", TEXT_SECONDARY, False),
    ("\u2022 Handles chunked/streaming output (partial lines)", TEXT_SECONDARY, False),
    ("\u2022 Deduplication via content hash set", TEXT_SECONDARY, False),
    ("\u2022 Mute system prevents echo detection:", TEXT_SECONDARY, False),
    ("    \u2192 mute() extend-only (never shortens)", GREEN_MID, False),
    ("    \u2192 unmute() force-cancels + snapshots grid", GREEN_MID, False),
    ("\u2022 Grid snapshot on unmute: marks visible protocol", TEXT_SECONDARY, False),
    ("  text as already-processed (no false re-detection)", TEXT_SECONDARY, False),
    ("\u2022 Injection grace period for freshly-injected protocols", TEXT_SECONDARY, False),
]
add_multiline(slide, Inches(1.1), Inches(3.1), Inches(5.2), Inches(3.8), scanner_lines, font_size=14)

# Right column — Orchestrator
add_rect(slide, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.8), BG_CARD, GREEN_DIM)
add_text(slide, Inches(7.2), Inches(2.5), Inches(5.2), Inches(0.4),
         "Orchestrator", font_size=18, color=GREEN_BRIGHT, bold=True)

orch_lines = [
    ("\u2022 Per-panel task queue with priority ordering", TEXT_SECONDARY, False),
    ("\u2022 Atomic message delivery via bracketed paste:", TEXT_SECONDARY, False),
    ("    \\x1b[200~ + text + \\x1b[201~ + \\r", GREEN_MID, False),
    ("\u2022 Chunked send for large messages (1024 byte chunks)", TEXT_SECONDARY, False),
    ("\u2022 15ms inter-chunk delay for PTY stability", TEXT_SECONDARY, False),
    ("\u2022 Mute scanner during send (5s ceiling)", TEXT_SECONDARY, False),
    ("\u2022 Unmute + grid snapshot after delivery", TEXT_SECONDARY, False),
    ("\u2022 Fire-and-forget ACK mute (1500ms, no timer)", TEXT_SECONDARY, False),
    ("\u2022 try/finally in processQueue prevents stalls", TEXT_SECONDARY, False),
    ("\u2022 ~600ms per message delivery (was 2300ms)", TEXT_SECONDARY, False),
]
add_multiline(slide, Inches(7.2), Inches(3.1), Inches(5.2), Inches(3.8), orch_lines, font_size=14)

slide_number(slide, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8 — Features
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "FEATURES", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "Everything you need. Nothing you don't.", font_size=38,
         color=GREEN_BRIGHT, bold=True)

features = [
    ("Multi-Agent Terminal", "Up to 4 agents in split panels.\nFull PTY with xterm-256color.", "Ctrl+2/3/4"),
    ("Inter-Agent Comms", "5 protocol commands: SEND, REPLY,\nBROADCAST, STATUS, QUERY.", "Commander Protocol"),
    ("120 Prompt Templates", "Curated library across 14 categories.\nBrowse and inject into any agent.", "Ctrl+B"),
    ("Dual-Panel File Manager", "Copy, move, delete, mkdir.\nNavigate without leaving.", "F6/F7/F8/F9"),
    ("10 Agent Types", "Claude, Codex, Gemini, Aider,\nCline, OpenCode, Goose, Kiro, Amp.", "F2"),
    ("Full Terminal Emulation", "xterm-256color, ANSI, mouse events.\nReal terminal, not a dumb pipe.", "FORCE_COLOR=1"),
]

for i, (title, desc, tag) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.8 + col * 4.0)
    top = Inches(2.3 + row * 2.4)
    add_rect(slide, left, top, Inches(3.6), Inches(2.1), BG_CARD, GREEN_DIM)
    add_text(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.0), Inches(0.4),
             title, font_size=17, color=TEXT_PRIMARY, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(0.7), Inches(3.0), Inches(1.0),
             desc, font_size=13, color=TEXT_SECONDARY)
    add_text(slide, left + Inches(0.3), top + Inches(1.7), Inches(3.0), Inches(0.3),
             tag, font_size=11, color=GREEN_MUTED)

slide_number(slide, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9 — Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "COMPARISON", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "Built different", font_size=38, color=GREEN_BRIGHT, bold=True)

# Table header
headers = ["", "Commander Protocol", "MCP (Anthropic)", "LangGraph / CrewAI"]
col_widths = [Inches(2.2), Inches(3.2), Inches(3.0), Inches(3.0)]
col_starts = [Inches(0.8)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

top = Inches(2.2)
row_h = Inches(0.42)

# Header bg
add_rect(slide, Inches(0.8), top, Inches(11.4), row_h, BG_CARD, GREEN_DIM)
for j, hdr in enumerate(headers):
    color = GREEN_BRIGHT if j == 1 else GREEN_MUTED
    add_text(slide, col_starts[j], top, col_widths[j], row_h,
             hdr, font_size=13, color=color, bold=True)

rows = [
    ["Architecture", "Terminal-native, text-based", "Client-server, JSON-RPC", "Framework, Python/JS SDK"],
    ["Integration Effort", "\u2713 Zero \u2014 just run agents", "SDK + server implementation", "Custom code + config"],
    ["Agent Support", "\u2713 Any CLI agent", "MCP-compatible only", "Framework-wrapped only"],
    ["Infrastructure", "\u2713 None \u2014 runs locally", "Server process required", "Runtime + dependencies"],
    ["API Keys Required", "\u2713 No (agents handle own)", "Per-server configuration", "Yes, centrally managed"],
    ["Real Terminal UI", "\u2713 Full TUI with panels", "\u2717 No UI", "\u2717 No UI"],
    ["File Management", "\u2713 Built-in dual-panel", "\u2717 Not included", "\u2717 Not included"],
    ["Learning Curve", "\u2713 Instant \u2014 text markers", "Moderate \u2014 protocol spec", "Steep \u2014 framework concepts"],
]

for i, row in enumerate(rows):
    y = top + row_h * (i + 1)
    bg = BG_CARD if i % 2 == 0 else BG_BLACK
    add_rect(slide, Inches(0.8), y, Inches(11.4), row_h, bg)
    for j, cell in enumerate(row):
        if j == 0:
            c = TEXT_PRIMARY
        elif j == 1:
            c = GREEN_BRIGHT if "\u2713" in cell else TEXT_PRIMARY
        else:
            c = TEXT_SECONDARY
        add_text(slide, col_starts[j], y, col_widths[j], row_h,
                 cell, font_size=12, color=c)

slide_number(slide, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10 — Supported Agents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "SUPPORTED AGENTS", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "10 agents. Any CLI tool via Generic adapter.", font_size=38,
         color=GREEN_BRIGHT, bold=True)

agents = [
    ("Claude Code", "Anthropic", "claude"),
    ("Codex CLI", "OpenAI", "codex"),
    ("Gemini CLI", "Google", "gemini"),
    ("Aider", "Paul Gauthier", "aider"),
    ("Cline", "VS Code agent", "cline"),
    ("OpenCode", "Open source", "opencode"),
    ("Goose", "Block", "goose"),
    ("Kiro", "AWS", "kiro"),
    ("Amp", "Sourcegraph", "amp"),
    ("Generic", "Any CLI tool", "custom command"),
]

for i, (name, vendor, cmd) in enumerate(agents):
    col = i % 5
    row = i // 5
    left = Inches(0.8 + col * 2.4)
    top = Inches(2.5 + row * 2.4)
    add_rect(slide, left, top, Inches(2.1), Inches(2.0), BG_CARD, GREEN_DIM)
    add_text(slide, left + Inches(0.2), top + Inches(0.3), Inches(1.7), Inches(0.4),
             name, font_size=16, color=GREEN_BRIGHT, bold=True)
    add_text(slide, left + Inches(0.2), top + Inches(0.8), Inches(1.7), Inches(0.3),
             vendor, font_size=12, color=TEXT_SECONDARY)
    add_text(slide, left + Inches(0.2), top + Inches(1.3), Inches(1.7), Inches(0.3),
             f"$ {cmd}", font_size=11, color=GREEN_MUTED)

slide_number(slide, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 11 — Keyboard Shortcuts
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
         "KEYBOARD SHORTCUTS", font_size=14, color=GREEN_MID, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
         "Function keys + Ctrl shortcuts", font_size=38, color=GREEN_BRIGHT, bold=True)

# F-keys column
add_rect(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(4.8), BG_CARD, GREEN_DIM)
add_text(slide, Inches(1.1), Inches(2.5), Inches(5.2), Inches(0.4),
         "Function Keys", font_size=18, color=GREEN_BRIGHT, bold=True)

fkeys = [
    ("F1", "Help"), ("F2", "Launch Agent"), ("F3", "Add Panel"),
    ("F4", "View File"), ("F5", "Edit File"), ("F6", "Copy"),
    ("F7", "Move / Rename"), ("F8", "Mkdir"), ("F9", "Delete"), ("F10", "Quit"),
]

fkey_lines = [(f"  {k:<6}  {v}", TEXT_SECONDARY if i > 0 else TEXT_SECONDARY, False)
              for i, (k, v) in enumerate(fkeys)]
add_multiline(slide, Inches(1.1), Inches(3.1), Inches(5.2), Inches(3.8), fkey_lines, font_size=14)

# Agent shortcuts column
add_rect(slide, Inches(6.9), Inches(2.3), Inches(5.8), Inches(4.8), BG_CARD, GREEN_DIM)
add_text(slide, Inches(7.2), Inches(2.5), Inches(5.2), Inches(0.4),
         "Agent Management", font_size=18, color=GREEN_BRIGHT, bold=True)

shortcuts = [
    ("Ctrl+O", "Send task to any agent"),
    ("Ctrl+P", "Inject protocol"),
    ("Ctrl+B", "Browse 120 prompt templates"),
    ("Ctrl+T", "Toggle file \u2194 terminal"),
    ("Ctrl+K", "Kill running session"),
    ("Ctrl+C", "Send interrupt to agent"),
    ("Ctrl+D", "Send EOF to agent"),
    ("Ctrl+2/3/4", "Switch panel layout"),
    ("Ctrl+W", "Remove active panel"),
    ("F12", "Inter-agent comm guide"),
]

shortcut_lines = [(f"  {k:<12}  {v}", TEXT_SECONDARY, False) for k, v in shortcuts]
add_multiline(slide, Inches(7.2), Inches(3.1), Inches(5.2), Inches(3.8), shortcut_lines, font_size=14)

slide_number(slide, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 12 — Get Started
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(5.667), Inches(0.5), Inches(2.0))

add_text(slide, Inches(0), Inches(2.7), W, Inches(0.8),
         "Get Started", font_size=44, color=GREEN_BRIGHT, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(3.6), W, Inches(0.5),
         "One command to install. One command to start collaborating.",
         font_size=20, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Install box
add_rect(slide, Inches(3.5), Inches(4.4), Inches(6.3), Inches(1.8), BG_CARD, GREEN_DIM)
install_lines = [
    ("$ npm install -g agents-commander", GREEN_BRIGHT, True),
    ("", TEXT_PRIMARY, False),
    ("$ cd your-project/", TEXT_SECONDARY, False),
    ("$ agents-commander", GREEN_BRIGHT, True),
]
add_multiline(slide, Inches(3.8), Inches(4.6), Inches(5.7), Inches(1.5), install_lines, font_size=16)

add_text(slide, Inches(0), Inches(6.5), W, Inches(0.4),
         "github.com/lkalinowski/agents-commander",
         font_size=16, color=GREEN_MID, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(7.0), W, Inches(0.4),
         "Node.js 18+  \u2022  macOS / Linux  \u2022  CC BY-NC 4.0",
         font_size=13, color=GREEN_MUTED, alignment=PP_ALIGN.CENTER)

slide_number(slide, 12, TOTAL_SLIDES)

# ── Save ──
prs.save(OUTPUT_PATH)
print(f"Presentation saved to: {OUTPUT_PATH}")
print(f"  {TOTAL_SLIDES} slides, green military theme")
